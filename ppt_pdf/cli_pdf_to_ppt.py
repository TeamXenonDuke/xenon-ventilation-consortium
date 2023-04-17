import io
import os

from absl import app, flags
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader, PdfWriter

# define flags
FLAGS = flags.FLAGS

flags.DEFINE_string("pdffile", None, "input pdf file path.")

flags.mark_flag_as_required("pdffile")


class PdfToPpt(object):
    """Converter of pdf file to pptx file.

    Attributes:
        pdf_file: str pdf file path
        ppt_file: str pptx output file path
    """

    def __init__(self, pdf_file: str):
        """Init PdfToPpt.

        Args:
            pdf_file: str path of pdf file.
        """
        self.total_pages = 1
        self.pdf_file = pdf_file
        self.ppt_file = pdf_file.replace(".pdf", ".pptx")

    def check_file_exist(self, file_path: str):
        """Check that file exists."""
        # self.log.info('Checking file - %s ' % file_path)
        if os.path.isfile(file_path):
            return True
        else:
            return False

    def pdf_to_image(self, pdf_file: str):
        """Convert pdf to image.

        Args:
            pdf_file: str pdf file path.
        """
        if not self.check_file_exist(pdf_file):
            return False
        image_file = pdf_file.replace(".pdf", ".jpg")
        try:
            pages = convert_from_path(pdf_file, 200)
            pages[0].save(image_file, "JPEG")
            return True
        except Exception:
            print(
                "Warning: Failed to convert .pdf to .jpg. Check if poppler is installed"
            )
            return False

    def pdf_splitter(self):
        """Split the pdfs into separate pages."""
        input_pdf = PdfReader(io.FileIO(self.pdf_file, "rb"))

        self.total_pages = len(input_pdf.pages)

        for page_number in range(self.total_pages):

            output = PdfWriter()
            output.add_page(input_pdf.pages[page_number])
            # new filename
            new_pdf = "_%s%s" % (str(page_number + 1), ".pdf")
            new_pdf = self.pdf_file.replace(".pdf", new_pdf)
            file_stream = io.FileIO(new_pdf, "wb")
            output.write(file_stream)
            file_stream.close()

            # calling pdf to image conversion
            self.pdf_to_image(new_pdf)

    def create_ppt(self):
        """Create the pptx file."""
        prs = Presentation()
        prs.slide_height = Inches(9)
        prs.slide_width = Inches(16)
        try:
            for slide_number in range(self.total_pages):
                img_path = self.pdf_file.replace(
                    ".pdf", "_%s%s" % (str(slide_number + 1), ".jpg")
                )
                # self.log.debug('%s' % img_path)
                new_slide = prs.slide_layouts[0]
                slide = prs.slides.add_slide(new_slide)
                subtitle = slide.placeholders[1]
                title = slide.shapes.title
                title.text = "Image %s " % str(slide_number + 1)
                left = Inches(0.3)
                top = Inches(0.5)
                height = Inches(7.8)
                pic = slide.shapes.add_picture(img_path, left, top, height=height)
                prs.save(self.ppt_file)
        except IOError:
            print("Warning: Error Creating PowerPoint")
            pass

    def execute(self):
        """Execute the ppt creation."""
        self.pdf_splitter()
        self.create_ppt()


def main(argv):
    """Command line interface."""
    filepath = FLAGS.pdffile
    PdfToPpt(pdf_file=filepath).execute()


if __name__ == "__main__":
    app.run(main)
